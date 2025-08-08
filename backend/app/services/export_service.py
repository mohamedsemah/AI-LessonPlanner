import io
import logging
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
import markdown
import re

logger = logging.getLogger(__name__)

class ExportService:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    async def export_to_powerpoint(self, course_content: Dict[str, Any], lesson_data: Dict[str, Any]) -> io.BytesIO:
        """Export course content to PowerPoint format"""
        try:
            self.logger.info("Creating PowerPoint presentation")
            
            # Create a new presentation
            prs = Presentation()
            
            # Get slide layouts
            title_slide_layout = prs.slide_layouts[0]  # Title slide
            content_slide_layout = prs.slide_layouts[1]  # Title and content
            section_header_layout = prs.slide_layouts[2]  # Section header
            
            # Add title slide
            title_slide = prs.slides.add_slide(title_slide_layout)
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            lesson_info = lesson_data.get("lesson_info", {})
            title.text = f"{lesson_info.get('course_title', 'Course')} - {lesson_info.get('lesson_topic', 'Lesson')}"
            subtitle.text = f"UDL-Compliant Course Content\n{len(course_content.get('slides', []))} slides • {course_content.get('estimated_duration', 0)} minutes"
            
            # Add slides for each Gagne event
            slides = course_content.get('slides', [])
            current_event = None
            
            for slide_data in slides:
                # Check if this is a new event
                if slide_data.get('gagne_event') != current_event:
                    current_event = slide_data.get('gagne_event')
                    
                    # Add section header for new event
                    section_slide = prs.slides.add_slide(section_header_layout)
                    section_title = section_slide.shapes.title
                    section_title.text = f"Event {current_event}: {slide_data.get('gagne_event_name', 'Unknown')}"
                
                # Add content slide
                content_slide = prs.slides.add_slide(content_slide_layout)
                slide_title = content_slide.shapes.title
                content_placeholder = content_slide.placeholders[1]
                
                # Set slide title
                slide_title.text = slide_data.get('title', 'Untitled Slide')
                
                # Add content
                content_text = ""
                
                # Add main content
                if slide_data.get('content'):
                    content_text += f"{slide_data['content']}\n\n"
                
                # Add visual elements description
                visual_elements = slide_data.get('visual_elements', [])
                if visual_elements:
                    content_text += "Visual Elements:\n"
                    for element in visual_elements:
                        if isinstance(element, dict):
                            content_text += f"• {element.get('description', 'Visual element')}\n"
                        else:
                            content_text += f"• {element}\n"
                    content_text += "\n"
                
                # Add audio script
                if slide_data.get('audio_script'):
                    content_text += f"Audio Script: {slide_data['audio_script']}\n\n"
                
                # Add speaker notes
                if slide_data.get('speaker_notes'):
                    content_text += f"Speaker Notes: {slide_data['speaker_notes']}\n\n"
                
                # Add UDL guidelines
                udl_guidelines = slide_data.get('udl_guidelines', [])
                if udl_guidelines:
                    content_text += "UDL Guidelines:\n"
                    for guideline in udl_guidelines:
                        content_text += f"• {guideline}\n"
                
                content_placeholder.text = content_text
            
            # Save to buffer
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            self.logger.info(f"PowerPoint export completed with {len(prs.slides)} slides")
            return buffer
            
        except Exception as e:
            self.logger.error(f"Error in PowerPoint export: {str(e)}")
            raise Exception(f"Failed to create PowerPoint: {str(e)}")
    
    async def export_to_pdf(self, course_content: Dict[str, Any], lesson_data: Dict[str, Any]) -> io.BytesIO:
        """Export course content to PDF format"""
        try:
            self.logger.info("Creating PDF document")
            
            # Create PDF buffer
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
            
            # Get styles
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=16,
                spaceAfter=30,
                alignment=TA_CENTER
            )
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontSize=14,
                spaceAfter=12,
                spaceBefore=20
            )
            normal_style = styles['Normal']
            
            # Build story (content)
            story = []
            
            # Add title page
            lesson_info = lesson_data.get("lesson_info", {})
            title_text = f"{lesson_info.get('course_title', 'Course')} - {lesson_info.get('lesson_topic', 'Lesson')}"
            story.append(Paragraph(title_text, title_style))
            story.append(Spacer(1, 20))
            
            # Add subtitle
            subtitle_text = f"UDL-Compliant Course Content"
            story.append(Paragraph(subtitle_text, heading_style))
            story.append(Spacer(1, 20))
            
            # Add metadata
            metadata = [
                f"Total Slides: {len(course_content.get('slides', []))}",
                f"Estimated Duration: {course_content.get('estimated_duration', 0)} minutes",
                f"UDL Compliance: {course_content.get('udl_compliance_report', {}).get('overall_compliance', 0) * 100:.1f}%"
            ]
            
            for meta in metadata:
                story.append(Paragraph(meta, normal_style))
                story.append(Spacer(1, 6))
            
            story.append(PageBreak())
            
            # Add slides content
            slides = course_content.get('slides', [])
            current_event = None
            
            for i, slide_data in enumerate(slides, 1):
                # Check if this is a new event
                if slide_data.get('gagne_event') != current_event:
                    current_event = slide_data.get('gagne_event')
                    
                    # Add event header
                    event_title = f"Event {current_event}: {slide_data.get('gagne_event_name', 'Unknown')}"
                    story.append(Paragraph(event_title, heading_style))
                    story.append(Spacer(1, 12))
                
                # Add slide content
                slide_title = f"Slide {i}: {slide_data.get('title', 'Untitled')}"
                story.append(Paragraph(slide_title, heading_style))
                story.append(Spacer(1, 6))
                
                # Add main content
                if slide_data.get('content'):
                    # Convert markdown to plain text for PDF
                    content_text = self._convert_markdown_to_text(slide_data['content'])
                    story.append(Paragraph(content_text, normal_style))
                    story.append(Spacer(1, 6))
                
                # Add visual elements
                visual_elements = slide_data.get('visual_elements', [])
                if visual_elements:
                    story.append(Paragraph("Visual Elements:", normal_style))
                    for element in visual_elements:
                        if isinstance(element, dict):
                            element_text = f"• {element.get('description', 'Visual element')}"
                        else:
                            element_text = f"• {element}"
                        story.append(Paragraph(element_text, normal_style))
                    story.append(Spacer(1, 6))
                
                # Add audio script
                if slide_data.get('audio_script'):
                    story.append(Paragraph("Audio Script:", normal_style))
                    story.append(Paragraph(slide_data['audio_script'], normal_style))
                    story.append(Spacer(1, 6))
                
                # Add speaker notes
                if slide_data.get('speaker_notes'):
                    story.append(Paragraph("Speaker Notes:", normal_style))
                    story.append(Paragraph(slide_data['speaker_notes'], normal_style))
                    story.append(Spacer(1, 6))
                
                # Add UDL guidelines
                udl_guidelines = slide_data.get('udl_guidelines', [])
                if udl_guidelines:
                    story.append(Paragraph("UDL Guidelines:", normal_style))
                    for guideline in udl_guidelines:
                        story.append(Paragraph(f"• {guideline}", normal_style))
                
                story.append(Spacer(1, 20))
            
            # Build PDF
            doc.build(story)
            buffer.seek(0)
            
            self.logger.info(f"PDF export completed with {len(slides)} slides")
            return buffer
            
        except Exception as e:
            self.logger.error(f"Error in PDF export: {str(e)}")
            raise Exception(f"Failed to create PDF: {str(e)}")
    
    def _convert_markdown_to_text(self, markdown_text: str) -> str:
        """Convert markdown text to plain text for PDF"""
        if not markdown_text:
            return ""
        
        # Remove markdown formatting
        text = markdown_text
        
        # Remove headers
        text = re.sub(r'^#+\s*', '', text, flags=re.MULTILINE)
        
        # Remove bold and italic
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        
        # Remove code blocks
        text = re.sub(r'```.*?```', '', text, flags=re.DOTALL)
        text = re.sub(r'`(.*?)`', r'\1', text)
        
        # Remove links
        text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
        
        # Remove lists
        text = re.sub(r'^\s*[-*+]\s*', '• ', text, flags=re.MULTILINE)
        text = re.sub(r'^\s*\d+\.\s*', '', text, flags=re.MULTILINE)
        
        # Clean up extra whitespace
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = text.strip()
        
        return text 